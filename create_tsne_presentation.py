import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import io
from image_scatter import visualize_tsne_images

def save_plot_to_bytes(fig):
    """Convert matplotlib figure to bytes object"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight', dpi=300)
    buf.seek(0)
    return buf

def create_tsne_presentation(plot_data_list, output_file='tsne_visualization.pptx'):
    """
    Create a PowerPoint presentation with t-SNE plots
    
    Parameters:
    -----------
    plot_data_list : list of dictionaries
        Each dictionary should contain:
        - 'tsne_results': t-SNE coordinates
        - 'images': corresponding images
        - 'title': title for the slide (optional)
        - 'n_samples': number of samples to plot (optional)
        - 'zoom': zoom factor for images (optional)
    output_file : str
        Name of the output PowerPoint file
    """
    # Create presentation
    prs = Presentation()
    
    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "t-SNE Visualizations"
    
    # For each plot data, create a slide
    for i, plot_data in enumerate(plot_data_list, 1):
        # Extract parameters
        tsne_results = plot_data['tsne_results']
        images = plot_data['images']
        title = plot_data.get('title', f't-SNE Plot {i}')
        n_samples = plot_data.get('n_samples', 1000)
        zoom = plot_data.get('zoom', 0.3)
        
        # Create t-SNE plot
        fig = visualize_tsne_images(
            tsne_results=tsne_results,
            images=images,
            n_samples=n_samples,
            zoom=zoom
        )
        
        # Convert plot to bytes
        img_bytes = save_plot_to_bytes(fig)
        plt.close(fig)  # Close the figure to free memory
        
        # Add content slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout with title and content
        
        # Add title to slide
        slide.shapes.title.text = title
        
        # Add plot to slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)  # Adjust these values as needed
        height = Inches(5.5)  # Adjust these values as needed
        
        slide.shapes.add_picture(img_bytes, left, top, width, height)
    
    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example usage:
"""
# Assuming you have multiple t-SNE results and corresponding images:
plot_data_list = [
    {
        'tsne_results': tsne_results1,
        'images': images1,
        'title': 'Dataset 1 t-SNE Visualization',
        'n_samples': 800,
        'zoom': 0.3
    },
    {
        'tsne_results': tsne_results2,
        'images': images2,
        'title': 'Dataset 2 t-SNE Visualization',
        'n_samples': 1000,
        'zoom': 0.25
    }
]

create_tsne_presentation(plot_data_list, 'my_tsne_plots.pptx')
""" 